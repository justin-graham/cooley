"""
Document parsing utilities for PDF, DOCX, XLSX, and PPTX files.
Gracefully handles errors - returns error message instead of crashing.
"""

import os
import logging
import zipfile
import tempfile
import chardet
from typing import List, Dict, Any, Tuple
import pymupdf4llm
import pymupdf as fitz
from docx import Document
import pandas as pd
from pptx import Presentation

logger = logging.getLogger(__name__)

# Supported file extensions for document processing
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.xlsx', '.pptx', '.png', '.jpg', '.jpeg', '.gif', '.bmp'}

# File extensions that are always skipped (system files, etc.)
SKIP_EXTENSIONS = {'.ds_store', '.gitkeep', '.gitignore', '.thumbs.db'}

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


def unzip_file(zip_path: str) -> List[str]:
    """
    Extract a zip file to a temporary directory.

    Args:
        zip_path: Path to the .zip file

    Returns:
        List of extracted file paths
    """
    extract_dir = tempfile.mkdtemp()

    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)

    # Recursively find all files (not directories)
    file_paths = []
    for root, dirs, files in os.walk(extract_dir):
        # Skip __MACOSX directories entirely
        if '__MACOSX' in root:
            continue
        for file in files:
            # Skip hidden files and system files
            if not file.startswith('.'):
                file_paths.append(os.path.join(root, file))

    return file_paths


def extract_from_pdf(file_path: str) -> str:
    """
    Extract text from PDF using pymupdf4llm (outputs Markdown for structure).

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted text as Markdown string

    Raises:
        Exception if extraction fails
    """
    # pymupdf4llm returns markdown-formatted text, preserving tables and structure
    markdown_text = pymupdf4llm.to_markdown(file_path)
    # Remove NULL bytes and control characters that break PostgreSQL
    markdown_text = markdown_text.replace('\x00', '').replace('\r', '\n')
    return markdown_text


def extract_from_pdf_fallback(file_path: str) -> str:
    """
    Fast fallback PDF text extraction using basic PyMuPDF.
    Used when pymupdf4llm times out on complex documents.

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted plain text string

    Raises:
        Exception if extraction fails
    """
    doc = fitz.open(file_path)
    text_parts = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        # Remove NULL bytes and control characters that break PostgreSQL
        text = text.replace('\x00', '').replace('\r', '\n')
        if text.strip():
            text_parts.append(f"--- Page {page_num + 1} ---\n{text}")

    doc.close()
    return '\n\n'.join(text_parts)


def extract_from_docx(file_path: str) -> str:
    """
    Extract text from DOCX file.

    Args:
        file_path: Path to DOCX file

    Returns:
        Extracted text string

    Raises:
        Exception if extraction fails
    """
    doc = Document(file_path)

    # Extract all paragraphs
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

    # Also extract text from tables
    tables_text = []
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            tables_text.append(' | '.join(row_text))

    # Combine paragraphs and tables
    full_text = '\n'.join(paragraphs)
    if tables_text:
        full_text += '\n\n--- Tables ---\n' + '\n'.join(tables_text)

    return full_text


def extract_from_xlsx(file_path: str) -> str:
    """
    Extract text and data from Excel file.

    Args:
        file_path: Path to XLSX file

    Returns:
        Extracted text (sheet names + cell values as string)

    Raises:
        Exception if extraction fails
    """
    # Read all sheets
    excel_file = pd.ExcelFile(file_path)

    all_text = []
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(excel_file, sheet_name=sheet_name)

        # Convert dataframe to text representation
        sheet_text = f"Sheet: {sheet_name}\n"
        sheet_text += df.to_string(index=False)
        all_text.append(sheet_text)

    return '\n\n'.join(all_text)


def extract_from_pptx(file_path: str) -> str:
    """
    Extract text from PowerPoint file.

    Args:
        file_path: Path to PPTX file

    Returns:
        Extracted text from all slides

    Raises:
        Exception if extraction fails
    """
    prs = Presentation(file_path)

    all_text = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = [f"--- Slide {i} ---"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

        all_text.append('\n'.join(slide_text))

    return '\n\n'.join(all_text)


def extract_from_image(file_path: str) -> str:
    """
    Extract basic information from image files (PNG, JPG, etc.).

    Args:
        file_path: Path to image file

    Returns:
        Basic image metadata as text

    Raises:
        Exception if extraction fails
    """
    if not PIL_AVAILABLE:
        return "Image file (PIL not available for text extraction)"

    img = Image.open(file_path)
    metadata = f"Image: {img.format} format, {img.size[0]}x{img.size[1]} pixels"

    # Try to extract any embedded text from metadata
    if hasattr(img, 'info') and img.info:
        for key, value in img.info.items():
            if isinstance(value, str) and len(value) < 200:
                metadata += f"\n{key}: {value}"

    return metadata


def parse_pdf_with_bboxes(filepath: str) -> Dict[str, Any]:
    """
    Extract text + bounding boxes for each text span from PDF.
    Used for generating document preview screenshots with highlights.

    Args:
        filepath: Path to PDF file

    Returns:
        Dictionary with keys:
        - full_text: str (complete document text for AI extraction)
        - text_spans: List[Dict] (each span has: text, page, bbox, char_offset_start, char_offset_end, font_size)

    Example text_span:
        {
            'text': '1,000,000',
            'page': 1,
            'bbox': [x0, y0, x1, y1],  # PDF coordinates
            'char_offset_start': 245,
            'char_offset_end': 254,
            'font_size': 12.0
        }
    """
    doc = fitz.open(filepath)
    text_spans = []
    full_text = ""
    char_offset = 0

    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

        for block in blocks:
            if "lines" in block:
                for line in block["lines"]:
                    for span in line["spans"]:
                        span_text = span['text']
                        text_spans.append({
                            'text': span_text,
                            'page': page_num,
                            'bbox': span['bbox'],  # [x0, y0, x1, y1]
                            'char_offset_start': char_offset,
                            'char_offset_end': char_offset + len(span_text),
                            'font_size': span.get('size', 12.0)
                        })
                        full_text += span_text
                        char_offset += len(span_text)

    doc.close()

    return {
        'full_text': full_text,
        'text_spans': text_spans
    }


def convert_to_pdf(filepath: str, output_dir: str) -> str:
    """
    Convert DOCX/XLSX/PPTX to PDF using LibreOffice headless.
    Required for generating document previews from non-PDF files.

    Args:
        filepath: Path to DOCX/XLSX/PPTX file
        output_dir: Directory for output PDF

    Returns:
        Path to generated PDF file

    Raises:
        RuntimeError: If LibreOffice is not installed
        subprocess.CalledProcessError: If conversion fails
    """
    import subprocess

    # Check if libreoffice is available
    try:
        subprocess.run(
            ['libreoffice', '--version'],
            capture_output=True,
            check=True,
            timeout=5
        )
    except (subprocess.CalledProcessError, FileNotFoundError):
        raise RuntimeError(
            "LibreOffice not installed. Required for DOCX/XLSX/PPTX preview generation. "
            "Install with: brew install libreoffice (macOS) or apt-get install libreoffice (Linux)"
        )

    # Convert to PDF
    logger.info(f"Converting {os.path.basename(filepath)} to PDF for preview generation")
    subprocess.run([
        'libreoffice', '--headless', '--convert-to', 'pdf',
        '--outdir', output_dir,
        filepath
    ], check=True, timeout=30)

    # Return path to generated PDF
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    output_pdf = os.path.join(output_dir, f"{base_name}.pdf")

    if not os.path.exists(output_pdf):
        raise RuntimeError(f"PDF conversion failed - output file not created: {output_pdf}")

    return output_pdf


def parse_document_with_paragraphs(file_path: str) -> Dict[str, Any]:
    """
    Parse a document file and extract text with paragraph numbering.
    Each substantive paragraph gets a sequential number for citation purposes.

    Args:
        file_path: Path to document file

    Returns:
        Dictionary with keys:
        - filename: str
        - type: str (pdf, docx, xlsx, pptx, or unknown)
        - full_text: str (complete extracted text)
        - paragraphs: List[Dict] with {number: int, text: str}
        - error: str (error message if parsing failed, otherwise absent)
    """
    # First get the basic parsed result
    basic_result = parse_document(file_path)

    result = {
        'filename': basic_result['filename'],
        'type': basic_result['type'],
        'full_text': basic_result.get('text', ''),
        'paragraphs': []
    }

    # If there was an error in basic parsing, return it
    if 'error' in basic_result:
        result['error'] = basic_result['error']
        return result

    # Split text into paragraphs and number them
    full_text = basic_result.get('text', '')
    if not full_text:
        return result

    # Split on double newlines (paragraph breaks) or single newlines for tighter docs
    raw_paragraphs = full_text.split('\n')

    paragraph_number = 1
    current_paragraph = []

    for line in raw_paragraphs:
        stripped_line = line.strip()

        # Skip truly empty lines
        if not stripped_line:
            # If we have accumulated text, save it as a paragraph
            if current_paragraph:
                para_text = ' '.join(current_paragraph)
                if len(para_text) > 20:  # Minimum 20 chars to be substantive
                    result['paragraphs'].append({
                        'number': paragraph_number,
                        'text': para_text
                    })
                    paragraph_number += 1
                current_paragraph = []
            continue

        # Accumulate non-empty lines
        current_paragraph.append(stripped_line)

    # Don't forget the last paragraph
    if current_paragraph:
        para_text = ' '.join(current_paragraph)
        if len(para_text) > 20:
            result['paragraphs'].append({
                'number': paragraph_number,
                'text': para_text
            })

    return result


def parse_document(file_path: str) -> Dict[str, Any]:
    """
    Parse a document file and extract text content.
    Gracefully handles errors - returns error dict instead of raising.

    Args:
        file_path: Path to document file

    Returns:
        Dictionary with keys:
        - filename: str
        - type: str (pdf, docx, xlsx, pptx, or unknown)
        - text: str (extracted text, or empty if error)
        - error: str (error message if parsing failed, otherwise absent)
    """
    filename = os.path.basename(file_path)
    file_ext = os.path.splitext(filename)[1].lower()

    result = {
        'filename': filename,
        'type': file_ext.lstrip('.'),
        'text': ''
    }

    try:
        # Route to appropriate parser based on extension
        if file_ext == '.pdf':
            try:
                result['text'] = extract_from_pdf(file_path)
            except NameError as e:
                # Library bug in pymupdf4llm - use fallback extraction
                logger.warning(f"pymupdf4llm bug for {filename}, using fallback: {e}")
                result['text'] = extract_from_pdf_fallback(file_path)
        elif file_ext == '.docx':
            result['text'] = extract_from_docx(file_path)
        elif file_ext == '.xlsx':
            result['text'] = extract_from_xlsx(file_path)
        elif file_ext == '.pptx':
            result['text'] = extract_from_pptx(file_path)
        elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            result['text'] = extract_from_image(file_path)
        elif not file_ext or file_ext == '.':
            # Files without extension - try to read as text
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    result['text'] = f.read(1000)  # First 1000 chars
            except:
                result['error'] = f"Unsupported file type: no extension"
                return result
        else:
            # Unsupported file type
            result['error'] = f"Unsupported file type: {file_ext}"
            return result

        # Check if extraction yielded any content
        if not result['text'] or len(result['text'].strip()) < 10:
            result['error'] = "Document appears to be empty or unreadable"

    except Exception as e:
        result['error'] = f"Failed to parse: {str(e)}"

    return result


def validate_zip_file(file_path: str, max_size_mb: int = 50) -> tuple[bool, str]:
    """
    Validate that a file is a valid zip and within size limits.

    Args:
        file_path: Path to the file
        max_size_mb: Maximum allowed size in megabytes

    Returns:
        Tuple of (is_valid, error_message)
    """
    # Check if file exists
    if not os.path.exists(file_path):
        return False, "File does not exist"

    # Check file size
    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
    if file_size_mb > max_size_mb:
        return False, f"File too large ({file_size_mb:.1f}MB). Maximum is {max_size_mb}MB"

    # Check if it's a valid zip
    if not zipfile.is_zipfile(file_path):
        return False, "File is not a valid ZIP archive"

    # Check decompressed size and file count to prevent zip bombs
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            total_uncompressed = sum(info.file_size for info in zf.infolist())
            file_count = len([info for info in zf.infolist() if not info.is_dir()])
            max_uncompressed_mb = 500
            max_file_count = 2000
            if total_uncompressed > max_uncompressed_mb * 1024 * 1024:
                return False, f"Archive too large when extracted ({total_uncompressed / (1024*1024):.0f}MB). Maximum is {max_uncompressed_mb}MB"
            if file_count > max_file_count:
                return False, f"Archive contains too many files ({file_count}). Maximum is {max_file_count}"
    except zipfile.BadZipFile:
        return False, "File is a corrupt ZIP archive"

    return True, ""


def detect_encoding(file_path: str) -> str:
    """
    Detect the encoding of a text file using chardet.

    Args:
        file_path: Path to the file

    Returns:
        Detected encoding (e.g., 'utf-8', 'latin-1', 'windows-1252')
    """
    with open(file_path, 'rb') as f:
        raw_data = f.read(10000)  # Read first 10KB for detection
        result = chardet.detect(raw_data)
        return result.get('encoding', 'utf-8') or 'utf-8'


def normalize_text_encoding(text: str) -> str:
    """
    Normalize text by removing problematic characters and ensuring clean UTF-8.

    Args:
        text: Input text string

    Returns:
        Cleaned text safe for PostgreSQL storage
    """
    if not text:
        return ''
    # Remove NULL bytes and control characters that break PostgreSQL
    text = text.replace('\x00', '')
    # Normalize line endings
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    return text


def prevalidate_zip_contents(zip_path: str) -> Dict[str, Any]:
    """
    Pre-validate zip contents before extraction.
    Scans for potential issues without extracting files.

    Args:
        zip_path: Path to the zip file

    Returns:
        Dictionary with validation results:
        - valid_files: List of file paths that can be processed
        - skipped_files: List of {path, reason} for files that will be skipped
        - warnings: List of warning messages
        - has_nested_folders: Boolean
        - duplicate_basenames: Dict mapping basename to list of full paths
    """
    result = {
        'valid_files': [],
        'skipped_files': [],
        'warnings': [],
        'has_nested_folders': False,
        'duplicate_basenames': {}
    }

    seen_basenames = {}  # Track duplicate filenames

    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            for zip_info in zip_ref.infolist():
                # Skip directories
                if zip_info.is_dir():
                    continue

                file_path = zip_info.filename
                basename = os.path.basename(file_path)

                # Skip hidden files, system files, and anything inside __MACOSX
                if basename.startswith('.') or '__MACOSX' in file_path:
                    result['skipped_files'].append({
                        'path': file_path,
                        'reason': 'Hidden or system file'
                    })
                    continue

                # Check for nested folders
                if '/' in file_path:
                    result['has_nested_folders'] = True

                # Get file extension
                ext = os.path.splitext(basename)[1].lower()

                # Skip known system extensions
                if ext in SKIP_EXTENSIONS:
                    result['skipped_files'].append({
                        'path': file_path,
                        'reason': f'System file ({ext})'
                    })
                    continue

                # Check for supported extensions
                if ext and ext not in SUPPORTED_EXTENSIONS:
                    result['skipped_files'].append({
                        'path': file_path,
                        'reason': f'Unsupported file type ({ext})'
                    })
                    continue

                # Check for empty files
                if zip_info.file_size == 0:
                    result['skipped_files'].append({
                        'path': file_path,
                        'reason': 'Empty file (0 bytes)'
                    })
                    continue

                # Track duplicate basenames
                if basename in seen_basenames:
                    if basename not in result['duplicate_basenames']:
                        result['duplicate_basenames'][basename] = [seen_basenames[basename]]
                    result['duplicate_basenames'][basename].append(file_path)
                else:
                    seen_basenames[basename] = file_path

                # File passed all checks
                result['valid_files'].append(file_path)

    except zipfile.BadZipFile as e:
        result['warnings'].append(f"Corrupt zip file: {str(e)}")
    except Exception as e:
        result['warnings'].append(f"Error scanning zip: {str(e)}")

    # Generate warnings for duplicates
    for basename, paths in result['duplicate_basenames'].items():
        result['warnings'].append(
            f"Duplicate filename '{basename}' found in {len(paths)} locations - will use folder path to distinguish"
        )

    return result


def unzip_file_robust(zip_path: str) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    """
    Extract a zip file with robust handling of edge cases.
    Returns both successfully extracted files and skipped/failed files.

    Args:
        zip_path: Path to the .zip file

    Returns:
        Tuple of:
        - extracted_files: List of {path, original_name, parse_status}
        - skipped_files: List of {original_name, reason, parse_status}
    """
    extract_dir = tempfile.mkdtemp()
    extracted_files = []
    skipped_files = []

    # Pre-validate contents
    validation = prevalidate_zip_contents(zip_path)

    # Log warnings
    for warning in validation['warnings']:
        logger.warning(warning)

    # Track skipped files
    for skipped in validation['skipped_files']:
        skipped_files.append({
            'original_name': skipped['path'],
            'reason': skipped['reason'],
            'parse_status': 'skipped'
        })

    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            for file_path in validation['valid_files']:
                try:
                    # Handle duplicate filenames by including folder structure
                    basename = os.path.basename(file_path)
                    if basename in validation['duplicate_basenames']:
                        # Use folder path to create unique name
                        folder = os.path.dirname(file_path).replace('/', '_').replace('\\', '_')
                        unique_name = f"{folder}_{basename}" if folder else basename
                    else:
                        unique_name = basename

                    # Extract the file
                    extracted_path = zip_ref.extract(file_path, extract_dir)

                    # Verify file was extracted and is readable
                    if os.path.exists(extracted_path) and os.path.getsize(extracted_path) > 0:
                        extracted_files.append({
                            'path': extracted_path,
                            'original_name': unique_name,
                            'parse_status': 'pending'
                        })
                    else:
                        skipped_files.append({
                            'original_name': file_path,
                            'reason': 'File extraction failed or empty',
                            'parse_status': 'error'
                        })

                except Exception as e:
                    logger.warning(f"Failed to extract {file_path}: {e}")
                    skipped_files.append({
                        'original_name': file_path,
                        'reason': f'Extraction error: {str(e)}',
                        'parse_status': 'error'
                    })

    except Exception as e:
        logger.error(f"Failed to open zip file: {e}")
        raise

    return extracted_files, skipped_files


def parse_document_robust(file_path: str) -> Dict[str, Any]:
    """
    Parse a document with robust error handling and encoding detection.
    Returns parse_status to track success/failure.

    Args:
        file_path: Path to document file

    Returns:
        Dictionary with keys:
        - filename: str
        - type: str (pdf, docx, xlsx, pptx, or unknown)
        - text: str (extracted text, or empty if error)
        - parse_status: str ('success', 'partial', 'error')
        - error: str (error message if parsing failed, otherwise absent)
    """
    filename = os.path.basename(file_path)
    file_ext = os.path.splitext(filename)[1].lower()

    result = {
        'filename': filename,
        'type': file_ext.lstrip('.'),
        'text': '',
        'parse_status': 'pending'
    }

    try:
        # Route to appropriate parser based on extension
        if file_ext == '.pdf':
            try:
                result['text'] = extract_from_pdf(file_path)
                result['parse_status'] = 'success'
            except NameError as e:
                # Library bug in pymupdf4llm - use fallback extraction
                logger.warning(f"pymupdf4llm bug for {filename}, using fallback: {e}")
                result['text'] = extract_from_pdf_fallback(file_path)
                result['parse_status'] = 'partial'  # Fallback means partial extraction
            except Exception as e:
                # Try fallback for any PDF error
                logger.warning(f"PDF extraction failed for {filename}, trying fallback: {e}")
                try:
                    result['text'] = extract_from_pdf_fallback(file_path)
                    result['parse_status'] = 'partial'
                except Exception as fallback_error:
                    result['error'] = f"PDF parsing failed: {str(e)}"
                    result['parse_status'] = 'error'

        elif file_ext == '.docx':
            result['text'] = extract_from_docx(file_path)
            result['parse_status'] = 'success'

        elif file_ext == '.xlsx':
            result['text'] = extract_from_xlsx(file_path)
            result['parse_status'] = 'success'

        elif file_ext == '.pptx':
            result['text'] = extract_from_pptx(file_path)
            result['parse_status'] = 'success'

        elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            result['text'] = extract_from_image(file_path)
            result['parse_status'] = 'success'

        elif not file_ext or file_ext == '.':
            # Files without extension - try to read as text with encoding detection
            try:
                encoding = detect_encoding(file_path)
                with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                    result['text'] = f.read(10000)  # First 10000 chars
                result['parse_status'] = 'success'
            except Exception as e:
                result['error'] = f"Unsupported file type: no extension"
                result['parse_status'] = 'error'

        else:
            # Unsupported file type
            result['error'] = f"Unsupported file type: {file_ext}"
            result['parse_status'] = 'error'

        # Normalize text encoding for all successful parses
        if result['text']:
            result['text'] = normalize_text_encoding(result['text'])

        # Check if extraction yielded any content
        if result['parse_status'] != 'error' and (not result['text'] or len(result['text'].strip()) < 10):
            result['error'] = "Document appears to be empty or unreadable"
            result['parse_status'] = 'error'

    except Exception as e:
        result['error'] = f"Failed to parse: {str(e)}"
        result['parse_status'] = 'error'

    return result
